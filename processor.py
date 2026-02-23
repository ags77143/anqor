"""
processor.py — Extract text from YouTube, PDFs, PPTXs, TXT files.
"""
import io
import re


def extract_text_from_source(source_type: str, url: str = None, file=None) -> tuple:
    if source_type == "youtube":
        return _extract_youtube(url)
    elif source_type == "file":
        return _extract_file(file)
    else:
        raise ValueError(f"Unknown source type: {source_type}")


def _parse_youtube_id(url: str):
    url = url.strip()
    patterns = [
        r'(?:v=)([a-zA-Z0-9_-]{11})',
        r'(?:youtu\.be\/)([a-zA-Z0-9_-]{11})',
        r'(?:embed\/)([a-zA-Z0-9_-]{11})',
        r'(?:shorts\/)([a-zA-Z0-9_-]{11})',
        r'(?:live\/)([a-zA-Z0-9_-]{11})',
        r'^([a-zA-Z0-9_-]{11})$',
    ]
    for p in patterns:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


def _extract_youtube(url: str) -> tuple:
    from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled

    video_id = _parse_youtube_id(url)
    if not video_id:
        raise ValueError(f"Couldn't find a YouTube video ID in: {url}\nTry copying the URL directly from your browser.")

    try:
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
    except TranscriptsDisabled:
        raise ValueError(
            "This video has transcripts disabled by the uploader.\n\n"
            "**Workarounds:**\n"
            "- Go to [tactiq.io](https://tactiq.io) or [downsub.com](https://downsub.com), paste the YouTube URL, download the transcript as .txt, then upload it here\n"
            "- Or paste the transcript text directly using the 'Paste Transcript' option"
        )
    except Exception as e:
        raise ValueError(
            f"Couldn't access this video's transcripts.\n\n"
            "**Workarounds:**\n"
            "- Try [tactiq.io](https://tactiq.io) or [downsub.com](https://downsub.com) to get a transcript\n"
            "- Or paste the transcript directly\n\n"
            f"Technical error: {str(e)}"
        )

    # Try every possible transcript — manual first, then generated, then any language
    transcript = None
    errors = []

    # 1. Manual English
    try:
        transcript = transcript_list.find_manually_created_transcript(["en", "en-US", "en-GB", "en-AU"])
    except Exception as e:
        errors.append(f"manual EN: {e}")

    # 2. Auto-generated English
    if not transcript:
        try:
            transcript = transcript_list.find_generated_transcript(["en", "en-US", "en-GB", "en-AU"])
        except Exception as e:
            errors.append(f"auto EN: {e}")

    # 3. Any manual transcript (any language)
    if not transcript:
        try:
            for t in transcript_list:
                if not t.is_generated:
                    transcript = t
                    break
        except Exception as e:
            errors.append(f"any manual: {e}")

    # 4. Any transcript at all
    if not transcript:
        try:
            for t in transcript_list:
                transcript = t
                break
        except Exception as e:
            errors.append(f"any: {e}")

    if not transcript:
        raise ValueError(
            "Found the video but couldn't retrieve any transcript.\n\n"
            "**Workarounds:**\n"
            "- Use [tactiq.io](https://tactiq.io) — paste the YouTube URL and it exports the transcript\n"
            "- Use [downsub.com](https://downsub.com) to download subtitles as .txt then upload here\n"
            "- Switch to 'Paste Transcript' and paste it manually"
        )

    try:
        entries = transcript.fetch()
        text = " ".join(entry["text"] for entry in entries)
        text = re.sub(r"\[.*?\]", "", text)  # remove [Music] etc
        text = re.sub(r"\s+", " ", text).strip()
        if len(text) < 100:
            raise ValueError("Transcript was found but appears to be empty or too short.")
        return text, ""
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Failed to read transcript: {str(e)}")


def _extract_file(file) -> tuple:
    name = file.name.lower()
    if name.endswith(".txt"):
        return _extract_txt(file), file.name
    elif name.endswith(".pdf"):
        return _extract_pdf(file), file.name
    elif name.endswith(".pptx"):
        return _extract_pptx(file), file.name
    else:
        raise ValueError(f"Unsupported file type. Please upload PDF, PPTX, or TXT.")


def _extract_txt(file) -> str:
    content = file.read()
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1")


def _extract_pdf(file) -> str:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
        pages = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n\n".join(pages)
    except Exception as e:
        raise ValueError(f"Failed to read PDF: {str(e)}")


def _extract_pptx(file) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file.read()))
        slides = []
        for i, slide in enumerate(prs.slides):
            parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if parts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(parts))
        return "\n\n".join(slides)
    except Exception as e:
        raise ValueError(f"Failed to read PPTX: {str(e)}")
